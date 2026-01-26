#!/usr/bin/env python3
"""
Apply a PPTX plan.json (aura_plan_spec) to a PPTX package.

Produces a plan-aware apply_report.json with engine and touched_parts for Gate A.
"""
from __future__ import annotations

import argparse
import json
import shutil
import tempfile
import zipfile
from collections import Counter
from pathlib import Path
from typing import Any

import parts
import plan as planmod
import pptx_create
import pptx_ooxml


def _zip_dir(src_dir: Path, out_file: Path) -> None:
    out_file.parent.mkdir(parents=True, exist_ok=True)
    with zipfile.ZipFile(out_file, "w", compression=zipfile.ZIP_DEFLATED) as zf:
        for p in sorted(src_dir.rglob("*")):
            if p.is_dir():
                continue
            zf.write(p, p.relative_to(src_dir).as_posix())


def _has_external_links(pptx_path: Path) -> bool:
    import xml.etree.ElementTree as ET

    rel_ns = {"rels": "http://schemas.openxmlformats.org/package/2006/relationships"}
    with zipfile.ZipFile(pptx_path, "r") as zf:
        for name in zf.namelist():
            if not name.endswith(".rels"):
                continue
            try:
                root = ET.fromstring(zf.read(name))
            except Exception:
                continue
            for rel in root.findall("rels:Relationship", rel_ns):
                if rel.attrib.get("TargetMode") == "External":
                    return True
    return False


def _risk_summary_from_package(pptx_path: Path, before_parts: dict[str, Any]) -> dict[str, Any]:
    has_macros = False
    try:
        with zipfile.ZipFile(pptx_path, "r") as zf:
            has_macros = "ppt/vbaProject.bin" in set(zf.namelist())
    except Exception:
        has_macros = False

    has_external_links = _has_external_links(pptx_path)
    has_charts = bool(before_parts.get("risk_parts", {}).get("ppt/charts/"))
    has_controls = bool(before_parts.get("risk_parts", {}).get("ppt/activeX/") or before_parts.get("risk_parts", {}).get("ppt/controlProps/"))

    risk_level = "low"
    if has_macros or has_external_links:
        risk_level = "high"
    elif has_charts or has_controls:
        risk_level = "medium"

    return {
        "has_charts": has_charts,
        "has_pivots": False,
        "has_controls": has_controls,
        "has_macros": has_macros,
        "has_formulas": False,
        "has_external_links": has_external_links,
        "risk_level": risk_level,
    }


def _empty_parts_info() -> dict[str, Any]:
    return {"parts_count": 0, "risk_parts": {}, "parts_sample": []}


def _derive_touched_parts(pptx_path: Path) -> list[str]:
    touched: set[str] = {
        "[Content_Types].xml",
        "_rels/.rels",
        "ppt/presentation.xml",
        "ppt/_rels/presentation.xml.rels",
    }
    try:
        with zipfile.ZipFile(pptx_path, "r") as zf:
            for name in zf.namelist():
                if name.startswith("ppt/slides/slide") and name.endswith(".xml"):
                    touched.add(name)
                    rels = name.replace("ppt/slides/", "ppt/slides/_rels/") + ".rels"
                    if rels in zf.namelist():
                        touched.add(rels)
    except Exception:
        pass
    return sorted(touched)


def _apply_patch(in_pptx: Path, plan: dict[str, Any], out_pptx: Path) -> tuple[list[str], list[dict[str, Any]]]:
    tmpdir = Path(tempfile.mkdtemp(prefix="pptx_patch_"))
    touched: set[str] = set()
    op_results: list[dict[str, Any]] = []
    try:
        with zipfile.ZipFile(in_pptx, "r") as zf:
            zf.extractall(tmpdir)

        for op in plan.get("operations", []):
            if not isinstance(op, dict):
                continue
            k = op.get("op")
            if k == "replace_text":
                slide_id = int(op["slide_id"])
                shape_id = op.get("shape_id")
                find = op["find"]
                replace = op["replace"]

                slide_part = pptx_ooxml.get_slide_part(tmpdir, slide_id)
                slide_path = tmpdir / slide_part
                changed = pptx_ooxml.replace_text_in_slide(slide_path, find=find, replace=replace, shape_id=shape_id if isinstance(shape_id, str) else None)
                if changed:
                    touched.add(slide_part)
                op_results.append({"op": k, "slide_id": slide_id, "shape_id": shape_id or "", "replaced_runs": changed})

            elif k == "fill_placeholder":
                slide_id = int(op["slide_id"])
                placeholder_id = str(op["placeholder_id"])
                content = str(op["content"])

                slide_part = pptx_ooxml.get_slide_part(tmpdir, slide_id)
                slide_path = tmpdir / slide_part
                ok = pptx_ooxml.fill_placeholder(slide_path, placeholder_id=placeholder_id, content=content)
                if not ok:
                    raise ValueError(f"Placeholder not found on slide {slide_id}: {placeholder_id}")
                touched.add(slide_part)
                op_results.append({"op": k, "slide_id": slide_id, "placeholder_id": placeholder_id})

            elif k == "delete_slide":
                slide_id = int(op["slide_id"])
                info = pptx_ooxml.delete_slide(tmpdir, slide_id=slide_id)
                touched.update({"ppt/presentation.xml", "ppt/_rels/presentation.xml.rels"})
                op_results.append({"op": k, "slide_id": slide_id, "rid": info.get("rid", ""), "target": info.get("target", "")})

            elif k == "reorder_slides":
                order = [int(x) for x in op["order"]]
                pptx_ooxml.reorder_slides(tmpdir, order=order)
                touched.add("ppt/presentation.xml")
                op_results.append({"op": k, "order": order})

            elif k == "duplicate_slide":
                slide_id = int(op["slide_id"])
                after = int(op.get("after_slide_id", slide_id))
                info = pptx_ooxml.duplicate_slide(tmpdir, slide_id=slide_id, after_slide_id=after)
                touched.update({"ppt/presentation.xml", "ppt/_rels/presentation.xml.rels", "[Content_Types].xml"})
                new_part = info.get("new_part_uri")
                if isinstance(new_part, str) and new_part:
                    touched.add(new_part)
                    rels_part = new_part.replace("ppt/slides/", "ppt/slides/_rels/") + ".rels"
                    touched.add(rels_part)
                op_results.append({"op": k, **info})

            elif k == "add_slide":
                raise ValueError("add_slide is not supported in patch mode; use --mode python-pptx if available.")

        _zip_dir(tmpdir, out_pptx)
        return (sorted(touched), op_results)
    finally:
        shutil.rmtree(tmpdir, ignore_errors=True)


def _apply_python_pptx(in_pptx: Path, plan: dict[str, Any], out_pptx: Path) -> tuple[list[str], list[dict[str, Any]]]:
    try:
        from pptx import Presentation  # type: ignore
    except Exception as e:  # pragma: no cover
        raise ValueError(f"python-pptx not available in this environment: {e}") from e

    prs = Presentation(str(in_pptx))
    op_results: list[dict[str, Any]] = []

    def _slides_list():
        return list(prs.slides)

    def _remove_slide_at(idx0: int) -> None:
        sldIdLst = prs.slides._sldIdLst  # type: ignore[attr-defined]
        sldIdLst.remove(sldIdLst[idx0])

    def _reorder(order_1based: list[int]) -> None:
        sldIdLst = prs.slides._sldIdLst  # type: ignore[attr-defined]
        ids = list(sldIdLst)
        for el in ids:
            sldIdLst.remove(el)
        for idx in order_1based:
            sldIdLst.append(ids[idx - 1])

    for op in plan.get("operations", []):
        if not isinstance(op, dict):
            continue
        k = op.get("op")
        if k == "replace_text":
            slide_id = int(op["slide_id"])
            find = op["find"]
            replace = op["replace"]
            shape_id = op.get("shape_id")
            changed = 0
            slide = _slides_list()[slide_id - 1]
            for shape in slide.shapes:
                if not getattr(shape, "has_text_frame", False):
                    continue
                if isinstance(shape_id, str) and shape_id and str(getattr(shape, "shape_id", "")) != shape_id:
                    continue
                tf = shape.text_frame
                for p in tf.paragraphs:
                    for run in p.runs:
                        if find in run.text:
                            run.text = run.text.replace(find, replace)
                            changed += 1
            op_results.append({"op": k, "slide_id": slide_id, "shape_id": shape_id or "", "replaced_runs": changed})

        elif k == "fill_placeholder":
            slide_id = int(op["slide_id"])
            placeholder_id = str(op["placeholder_id"])
            content = str(op["content"])
            slide = _slides_list()[slide_id - 1]
            filled = 0
            for shape in slide.shapes:
                ph = getattr(shape, "placeholder_format", None)
                if ph is None:
                    continue
                # python-pptx placeholder idx is int; type can be inferred by id
                if placeholder_id.isdigit() and int(placeholder_id) != getattr(ph, "idx", -1):
                    continue
                if not placeholder_id.isdigit():
                    # best-effort type match via name
                    if placeholder_id.lower() not in str(getattr(shape, "name", "")).lower():
                        continue
                if getattr(shape, "has_text_frame", False):
                    shape.text = content
                    filled += 1
            if filled == 0:
                raise ValueError(f"Placeholder not found on slide {slide_id}: {placeholder_id}")
            op_results.append({"op": k, "slide_id": slide_id, "placeholder_id": placeholder_id, "filled": filled})

        elif k == "delete_slide":
            slide_id = int(op["slide_id"])
            _remove_slide_at(slide_id - 1)
            op_results.append({"op": k, "slide_id": slide_id})

        elif k == "reorder_slides":
            order = [int(x) for x in op["order"]]
            _reorder(order)
            op_results.append({"op": k, "order": order})

        elif k == "duplicate_slide":
            raise ValueError("duplicate_slide is only supported in patch mode.")

        elif k == "add_slide":
            after = int(op["after_slide_id"])
            layout_name = op["layout"]
            layout = None
            for l in prs.slide_layouts:
                if getattr(l, "name", "") == layout_name:
                    layout = l
                    break
            if layout is None:
                layout = prs.slide_layouts[0]
            prs.slides.add_slide(layout)
            # Insert after 'after' by moving last slide id element.
            if after > 0:
                sldIdLst = prs.slides._sldIdLst  # type: ignore[attr-defined]
                ids = list(sldIdLst)
                last = ids[-1]
                sldIdLst.remove(last)
                sldIdLst.insert(after, last)  # after is 1-based; insert at index=after
            op_results.append({"op": k, "after_slide_id": after, "layout": layout_name})

    out_pptx.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_pptx))

    # Best-effort touched parts: assume core presentation + slide XML.
    touched = ["ppt/presentation.xml", "ppt/_rels/presentation.xml.rels"]
    return (touched, op_results)


def apply_plan(in_pptx: Path | None, plan: dict[str, Any], out_pptx: Path, mode: str) -> dict[str, Any]:
    classification = planmod.classify_plan(plan)
    input_exists = in_pptx is not None and str(in_pptx) != "-" and in_pptx.exists()
    before = parts.list_parts(in_pptx) if input_exists else _empty_parts_info()

    chosen = mode
    if mode == "auto":
        if classification.get("is_create") or not input_exists:
            chosen = "create"
        else:
            chosen = "patch" if classification["is_patch_supported"] else "python-pptx"
    if chosen in {"patch", "python-pptx"} and not input_exists:
        raise ValueError("Input PPTX is required for patch/python-pptx modes. Use create mode or '-' input.")

    if chosen == "create":
        touched_parts, op_results = pptx_create.apply_create_plan(in_pptx if input_exists else None, plan, out_pptx)
        engine = "python-pptx-create"
        if not touched_parts:
            touched_parts = _derive_touched_parts(out_pptx)
    elif chosen == "patch":
        touched_parts, op_results = _apply_patch(in_pptx, plan, out_pptx)
        engine = "patch"
    else:
        touched_parts, op_results = _apply_python_pptx(in_pptx, plan, out_pptx)
        engine = "python-pptx"

    after = parts.list_parts(out_pptx)
    part_diff = parts.diff_parts(before, after)
    risk_summary = _risk_summary_from_package(in_pptx if input_exists else out_pptx, before if input_exists else after)

    op_kinds = [r.get("op") for r in op_results if isinstance(r, dict) and isinstance(r.get("op"), str)]
    counts = Counter(op_kinds)
    change_summary = ", ".join([f"{k}×{counts[k]}" for k in sorted(counts)]) if counts else ""

    return {
        "engine": engine,
        "touched_parts": touched_parts,
        "operations": op_results,
        "classification": classification,
        "parts_before": before,
        "parts_after": after,
        "parts_diff": part_diff,
        "change_summary": change_summary,
        "risk_summary": risk_summary,
    }


def main() -> None:
    ap = argparse.ArgumentParser(description="Apply aura_pptx plan.json to a PPTX file.")
    ap.add_argument("in_pptx", type=Path)
    ap.add_argument("plan_json", type=Path)
    ap.add_argument("out_pptx", type=Path)
    ap.add_argument("--mode", choices=["auto", "patch", "python-pptx", "create"], default="auto")
    ap.add_argument("--out", type=Path, default=None, help="Write apply_report.json")
    args = ap.parse_args()

    plan = planmod.load_and_validate_plan(args.plan_json)
    rep = apply_plan(args.in_pptx, plan, args.out_pptx, mode=args.mode)

    payload = json.dumps(rep, ensure_ascii=False, indent=2)
    if args.out:
        args.out.parent.mkdir(parents=True, exist_ok=True)
        args.out.write_text(payload, encoding="utf-8")
    else:
        print(payload)


if __name__ == "__main__":
    main()
