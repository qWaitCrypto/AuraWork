#!/usr/bin/env python3
"""
Create or extend PPTX decks using python-pptx.

This supports creation-oriented ops (create_deck, add_textbox, add_shape, etc.)
and a small subset of edit ops for mixed plans.
"""
from __future__ import annotations

from pathlib import Path
from typing import Any


def _load_pptx():
    try:
        from pptx import Presentation  # type: ignore
        from pptx.util import Inches, Pt  # type: ignore
        from pptx.dml.color import RGBColor  # type: ignore
        from pptx.enum.text import PP_ALIGN  # type: ignore
        from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE  # type: ignore
        return Presentation, Inches, Pt, RGBColor, PP_ALIGN, MSO_AUTO_SHAPE_TYPE
    except Exception as e:
        raise ValueError(f"python-pptx is required for create ops, but it is not available: {e}") from e


def _parse_color(value: str | None):
    if not value:
        return None
    cleaned = value.strip().lstrip("#")
    if len(cleaned) not in (6, 8):
        return None
    return cleaned[:6]


def _to_inches(value: Any) -> float:
    if isinstance(value, (int, float)):
        return float(value)
    raise ValueError(f"Expected numeric inches value, got {value!r}")


def _to_pt(value: Any) -> float | None:
    if value is None:
        return None
    if isinstance(value, (int, float)):
        return float(value)
    return None


def _find_layout(prs, *, name: str | None, idx: int | None):
    if idx is not None and 0 <= idx < len(prs.slide_layouts):
        return prs.slide_layouts[idx]
    if name:
        for layout in prs.slide_layouts:
            if getattr(layout, "name", "") == name:
                return layout
    return prs.slide_layouts[0]


def _get_slide(prs, slide_id: int):
    slides = list(prs.slides)
    if slide_id < 1 or slide_id > len(slides):
        raise ValueError(f"slide_id out of range: {slide_id}")
    return slides[slide_id - 1]


def _set_run_style(run, *, font_size: float | None, bold: bool | None, color: str | None):
    if font_size is not None:
        run.font.size = _load_pptx()[2](font_size)  # Pt
    if bold is not None:
        run.font.bold = bool(bold)
    hex_color = _parse_color(color) if color else None
    if hex_color:
        run.font.color.rgb = _load_pptx()[3].from_string(hex_color)


def _apply_alignment(p, align: str | None):
    if not align:
        return
    _, _, _, _, PP_ALIGN, _ = _load_pptx()
    mapping = {
        "left": PP_ALIGN.LEFT,
        "center": PP_ALIGN.CENTER,
        "right": PP_ALIGN.RIGHT,
        "justify": PP_ALIGN.JUSTIFY,
    }
    val = mapping.get(align.lower())
    if val is not None:
        p.alignment = val


def _add_textbox(slide, *, text: str, x: float, y: float, w: float, h: float, font_size: float | None, bold: bool | None, color: str | None, align: str | None):
    Presentation, Inches, Pt, RGBColor, PP_ALIGN, MSO_AUTO_SHAPE_TYPE = _load_pptx()
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    _set_run_style(run, font_size=font_size, bold=bold, color=color)
    _apply_alignment(p, align)
    return tb


def _set_fill(shape, color: str | None):
    hex_color = _parse_color(color) if color else None
    if not hex_color:
        return
    _, _, _, RGBColor, _, _ = _load_pptx()
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor.from_string(hex_color)


def _set_line(shape, color: str | None):
    hex_color = _parse_color(color) if color else None
    if not hex_color:
        return
    _, _, _, RGBColor, _, _ = _load_pptx()
    line = shape.line
    line.color.rgb = RGBColor.from_string(hex_color)


def _set_background(slide, color: str):
    hex_color = _parse_color(color)
    if not hex_color:
        return
    _, _, _, RGBColor, _, _ = _load_pptx()
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor.from_string(hex_color)


def _add_shape(slide, *, shape_name: str | None, x: float, y: float, w: float, h: float):
    Presentation, Inches, Pt, RGBColor, PP_ALIGN, MSO_AUTO_SHAPE_TYPE = _load_pptx()
    shape_map = {
        "rect": MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        "roundrect": MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        "round_rect": MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        "ellipse": MSO_AUTO_SHAPE_TYPE.OVAL,
        "triangle": MSO_AUTO_SHAPE_TYPE.ISOSCELES_TRIANGLE,
        "diamond": MSO_AUTO_SHAPE_TYPE.DIAMOND,
        "hex": MSO_AUTO_SHAPE_TYPE.HEXAGON,
    }
    key = (shape_name or "rect").replace(" ", "").lower()
    shape_type = shape_map.get(key, MSO_AUTO_SHAPE_TYPE.RECTANGLE)
    return slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))


def _set_notes(slide, text: str) -> None:
    notes = slide.notes_slide
    notes.notes_text_frame.text = text


def _reorder_slides(prs, order_1based: list[int]) -> None:
    sldIdLst = prs.slides._sldIdLst  # type: ignore[attr-defined]
    ids = list(sldIdLst)
    for el in ids:
        sldIdLst.remove(el)
    for idx in order_1based:
        sldIdLst.append(ids[idx - 1])


def _remove_slide_at(prs, idx0: int) -> None:
    sldIdLst = prs.slides._sldIdLst  # type: ignore[attr-defined]
    sldIdLst.remove(sldIdLst[idx0])


def _insert_after_last(prs, after_1based: int) -> None:
    if after_1based <= 0:
        return
    sldIdLst = prs.slides._sldIdLst  # type: ignore[attr-defined]
    ids = list(sldIdLst)
    last = ids[-1]
    sldIdLst.remove(last)
    sldIdLst.insert(after_1based, last)


def apply_create_plan(in_pptx: Path | None, plan: dict[str, Any], out_pptx: Path) -> tuple[list[str], list[dict[str, Any]]]:
    Presentation, Inches, Pt, RGBColor, PP_ALIGN, MSO_AUTO_SHAPE_TYPE = _load_pptx()

    is_new = in_pptx is None or str(in_pptx) == "-" or not in_pptx.exists()
    prs = Presentation() if is_new else Presentation(str(in_pptx))

    op_results: list[dict[str, Any]] = []

    for op in plan.get("operations", []):
        if not isinstance(op, dict):
            continue
        k = op.get("op")

        if k == "create_deck":
            if not is_new:
                raise ValueError("create_deck is only allowed when creating a new deck.")
            size = op.get("slide_size") or {}
            if isinstance(size, dict):
                w = size.get("width_in")
                h = size.get("height_in")
                if isinstance(w, (int, float)) and isinstance(h, (int, float)) and w > 0 and h > 0:
                    prs.slide_width = Inches(float(w))
                    prs.slide_height = Inches(float(h))
            op_results.append({"op": k})

        elif k == "set_slide_size":
            w = _to_inches(op.get("width_in"))
            h = _to_inches(op.get("height_in"))
            prs.slide_width = Inches(w)
            prs.slide_height = Inches(h)
            op_results.append({"op": k, "width_in": w, "height_in": h})

        elif k == "add_slide":
            layout_name = op.get("layout") if isinstance(op.get("layout"), str) else None
            layout_idx = op.get("layout_idx") if isinstance(op.get("layout_idx"), int) else None
            after = int(op.get("after_slide_id", 0))
            layout = _find_layout(prs, name=layout_name, idx=layout_idx)
            prs.slides.add_slide(layout)
            _insert_after_last(prs, after)
            op_results.append({"op": k, "after_slide_id": after, "layout": layout_name or "", "layout_idx": layout_idx})

        elif k == "add_title":
            slide_id = int(op["slide_id"])
            text = str(op.get("text", ""))
            slide = _get_slide(prs, slide_id)
            filled = False
            for shape in slide.shapes:
                ph = getattr(shape, "placeholder_format", None)
                if ph is None:
                    continue
                name = str(getattr(shape, "name", ""))
                if "title" not in name.lower():
                    continue
                if getattr(shape, "has_text_frame", False):
                    shape.text = text
                    filled = True
                    break
            if not filled:
                _add_textbox(
                    slide,
                    text=text,
                    x=0.7,
                    y=0.3,
                    w=12.0,
                    h=1.0,
                    font_size=_to_pt(op.get("font_size")) or 36,
                    bold=op.get("bold", True),
                    color=op.get("color"),
                    align=op.get("align"),
                )
            op_results.append({"op": k, "slide_id": slide_id})

        elif k == "add_textbox":
            slide_id = int(op["slide_id"])
            slide = _get_slide(prs, slide_id)
            _add_textbox(
                slide,
                text=str(op.get("text", "")),
                x=_to_inches(op.get("x")),
                y=_to_inches(op.get("y")),
                w=_to_inches(op.get("w")),
                h=_to_inches(op.get("h")),
                font_size=_to_pt(op.get("font_size")),
                bold=op.get("bold"),
                color=op.get("color"),
                align=op.get("align"),
            )
            op_results.append({"op": k, "slide_id": slide_id})

        elif k == "add_shape":
            slide_id = int(op["slide_id"])
            slide = _get_slide(prs, slide_id)
            shape = _add_shape(
                slide,
                shape_name=op.get("shape"),
                x=_to_inches(op.get("x")),
                y=_to_inches(op.get("y")),
                w=_to_inches(op.get("w")),
                h=_to_inches(op.get("h")),
            )
            _set_fill(shape, op.get("fill"))
            _set_line(shape, op.get("line"))
            text = op.get("text")
            if isinstance(text, str) and text:
                shape.text_frame.clear()
                p = shape.text_frame.paragraphs[0]
                run = p.add_run()
                run.text = text
                _set_run_style(run, font_size=_to_pt(op.get("font_size")), bold=op.get("bold"), color=op.get("color"))
                _apply_alignment(p, op.get("align"))
            op_results.append({"op": k, "slide_id": slide_id})

        elif k == "add_image":
            slide_id = int(op["slide_id"])
            slide = _get_slide(prs, slide_id)
            path = Path(str(op.get("path"))).expanduser().resolve()
            if not path.exists():
                raise ValueError(f"Image not found: {path}")
            slide.shapes.add_picture(
                str(path),
                Inches(_to_inches(op.get("x"))),
                Inches(_to_inches(op.get("y"))),
                Inches(_to_inches(op.get("w"))),
                Inches(_to_inches(op.get("h"))),
            )
            op_results.append({"op": k, "slide_id": slide_id, "path": str(path)})

        elif k == "set_slide_bg":
            slide_id = int(op["slide_id"])
            slide = _get_slide(prs, slide_id)
            _set_background(slide, str(op.get("color")))
            op_results.append({"op": k, "slide_id": slide_id})

        elif k == "add_notes":
            slide_id = int(op["slide_id"])
            slide = _get_slide(prs, slide_id)
            _set_notes(slide, str(op.get("text", "")))
            op_results.append({"op": k, "slide_id": slide_id})

        elif k == "replace_text":
            slide_id = int(op["slide_id"])
            find = op.get("find", "")
            replace = op.get("replace", "")
            shape_id = op.get("shape_id")
            changed = 0
            slide = _get_slide(prs, slide_id)
            for shape in slide.shapes:
                if not getattr(shape, "has_text_frame", False):
                    continue
                if isinstance(shape_id, str) and shape_id and str(getattr(shape, "shape_id", "")) != shape_id:
                    continue
                tf = shape.text_frame
                for p in tf.paragraphs:
                    for run in p.runs:
                        if find and find in run.text:
                            run.text = run.text.replace(find, replace)
                            changed += 1
            op_results.append({"op": k, "slide_id": slide_id, "shape_id": shape_id or "", "replaced_runs": changed})

        elif k == "fill_placeholder":
            slide_id = int(op["slide_id"])
            placeholder_id = str(op.get("placeholder_id", ""))
            content = str(op.get("content", ""))
            slide = _get_slide(prs, slide_id)
            filled = 0
            for shape in slide.shapes:
                ph = getattr(shape, "placeholder_format", None)
                if ph is None:
                    continue
                if placeholder_id.isdigit() and int(placeholder_id) != getattr(ph, "idx", -1):
                    continue
                if not placeholder_id.isdigit():
                    if placeholder_id.lower() not in str(getattr(shape, "name", "")).lower():
                        continue
                if getattr(shape, "has_text_frame", False):
                    shape.text = content
                    filled += 1
            if filled == 0:
                raise ValueError(f"Placeholder not found on slide {slide_id}: {placeholder_id}")
            op_results.append({"op": k, "slide_id": slide_id, "placeholder_id": placeholder_id, "filled": filled})

        elif k == "delete_slide":
            slide_id = int(op["slide_id"])
            _remove_slide_at(prs, slide_id - 1)
            op_results.append({"op": k, "slide_id": slide_id})

        elif k == "reorder_slides":
            order = [int(x) for x in op.get("order", [])]
            _reorder_slides(prs, order)
            op_results.append({"op": k, "order": order})

        elif k == "duplicate_slide":
            raise ValueError("duplicate_slide is not supported in create mode; use patch mode on existing decks.")

        else:
            raise ValueError(f"Unsupported op in create engine: {k}")

    out_pptx.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_pptx))

    return ([], op_results)


__all__ = ["apply_create_plan"]
